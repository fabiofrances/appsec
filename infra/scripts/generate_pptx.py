#!/usr/bin/env python3
"""Generate AppSec Palestra PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# ── Palette ──────────────────────────────────────────────────────────────────
BG     = RGBColor(0x0D, 0x1B, 0x2A)
BG2    = RGBColor(0x16, 0x27, 0x3E)
DGRAY  = RGBColor(0x2D, 0x3E, 0x50)
BLUE   = RGBColor(0x00, 0x8B, 0xD4)
GREEN  = RGBColor(0x00, 0xC8, 0x8A)
RED    = RGBColor(0xFF, 0x4B, 0x4B)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
YELLOW = RGBColor(0xFF, 0xD0, 0x00)
TEAL   = RGBColor(0x00, 0xBC, 0xD4)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xAA, 0xBB, 0xCC)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
blank = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def box(s, x, y, w, h, color, border=False):
    sh = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if not border:
        sh.line.fill.background()
    return sh


def txt(s, text, x, y, w, h, size=14, color=WHITE, bold=False,
        align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = "Segoe UI"
    return tb


def header(s, title, subtitle=None, accent=BLUE):
    box(s, 0, 0, 13.33, 0.07, accent)
    txt(s, title, 0.4, 0.12, 12.5, 0.75, size=30, bold=True)
    if subtitle:
        txt(s, subtitle, 0.4, 0.82, 12.5, 0.38, size=15, color=LGRAY, italic=True)
    box(s, 0.4, 1.18, 12.5, 0.025, accent)


def card(s, x, y, w, h, color, title, body_lines, icon=""):
    box(s, x, y, w, h, BG2)
    box(s, x, y, 0.06, h, color)
    txt(s, f"{icon}  {title}" if icon else title,
        x + 0.15, y + 0.08, w - 0.25, 0.45, size=15, bold=True, color=WHITE)
    for i, line in enumerate(body_lines):
        txt(s, line, x + 0.15, y + 0.58 + i * 0.42,
            w - 0.25, 0.38, size=12, color=LGRAY)


def pill(s, text, x, y, bg_color, fg=WHITE, size=12):
    w = max(len(text) * 0.115 + 0.35, 1.0)
    b = box(s, x, y, w, 0.33, bg_color)
    t = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.33))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = fg
    r.font.bold = True
    r.font.name = "Segoe UI"


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 1 — Título
# ═══════════════════════════════════════════════════════════════════════════════
s1 = slide()
box(s1, 0, 0, 3.8, 7.5, BG2)
box(s1, 3.72, 0, 0.09, 7.5, BLUE)

txt(s1, "🔒", 0.0, 1.3, 3.8, 2.5, size=110, align=PP_ALIGN.CENTER)

txt(s1, "Application Security", 4.0, 1.1, 9.0, 0.9, size=38, bold=True)
txt(s1, "Integrando Segurança ao", 4.0, 2.05, 9.0, 0.65, size=28, color=BLUE)
txt(s1, "Ciclo de Desenvolvimento", 4.0, 2.7, 9.0, 0.65, size=28, color=BLUE)
txt(s1, "Calculadora de IMC como caso de estudo", 4.0, 3.55, 9.0, 0.5, size=16, color=LGRAY, italic=True)

px = 4.0
for tech, color in [("NestJS", RED), ("Vue.js", GREEN), ("Docker", BLUE),
                    ("Trivy", RGBColor(0x1A, 0x1A, 0xCC)), ("OWASP ZAP", RED),
                    ("DefectDojo", PURPLE)]:
    pill(s1, tech, px, 4.3, color)
    px += len(tech) * 0.115 + 0.55

box(s1, 0, 7.1, 13.33, 0.4, BLUE)
txt(s1, "github.com/questores/appsec-bmi", 4.0, 7.12, 9.0, 0.35,
    size=13, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 2 — Agenda
# ═══════════════════════════════════════════════════════════════════════════════
s2 = slide()
header(s2, "Agenda", "O que vamos ver hoje", BLUE)

steps = [
    ("01", "Backend NestJS", "API REST com Swagger, OpenTelemetry e Prometheus", BLUE),
    ("02", "Frontend Vue.js", "SPA com OpenTelemetry Web e proxy nginx", GREEN),
    ("03", "Testes", "Jest (backend) + Vitest (frontend) — cobertura completa", ORANGE),
    ("04", "AppSec", "Trivy · OWASP Dep. Check · OWASP ZAP · DefectDojo", RED),
]
for i, (num, title, desc, color) in enumerate(steps):
    y = 1.45 + i * 1.42
    box(s2, 0.4, y, 0.65, 1.2, color)
    txt(s2, num, 0.4, y + 0.15, 0.65, 0.9, size=26, bold=True,
        align=PP_ALIGN.CENTER)
    box(s2, 1.15, y, 11.8, 1.2, BG2)
    txt(s2, title, 1.35, y + 0.08, 10.5, 0.5, size=20, bold=True)
    txt(s2, desc, 1.35, y + 0.62, 10.5, 0.45, size=14, color=LGRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 3 — O que é AppSec / Shift Left
# ═══════════════════════════════════════════════════════════════════════════════
s3 = slide()
header(s3, "O que é AppSec?", "Shift Left Security — segurança desde o início", RED)

txt(s3, "🔄  Shift Left Security", 0.5, 1.35, 12.0, 0.6, size=22, bold=True, color=YELLOW)
txt(s3, "Integrar segurança desde o início do desenvolvimento — não apenas ao final,",
    0.5, 1.95, 12.5, 0.45, size=16, color=WHITE)
txt(s3, "antes de ir para produção. Custo de correção: até 6× maior em produção.",
    0.5, 2.38, 12.5, 0.45, size=16, color=LGRAY)

stages = [
    ("💡", "Plan",   BG2),
    ("💻", "Code",   BLUE),
    ("🔨", "Build",  BLUE),
    ("🧪", "Test",   GREEN),
    ("🚀", "Deploy", ORANGE),
    ("🔍", "Monitor",PURPLE),
]
sx = 0.35
for icon, label, color in stages:
    box(s3, sx, 3.1, 2.0, 1.5, color)
    txt(s3, icon,  sx, 3.15, 2.0, 0.7, size=34, align=PP_ALIGN.CENTER)
    txt(s3, label, sx, 3.85, 2.0, 0.55, size=14, bold=True,
        align=PP_ALIGN.CENTER)
    if sx < 11.5:
        txt(s3, "→", sx + 1.95, 3.7, 0.35, 0.4, size=20, color=LGRAY,
            align=PP_ALIGN.CENTER)
    sx += 2.15

box(s3, 0.4, 4.85, 12.5, 0.65, DGRAY)
txt(s3, "🛡️  Neste projeto:  Trivy  ·  OWASP Dependency Check  ·  OWASP ZAP  ·  DefectDojo",
    0.7, 4.95, 12.0, 0.45, size=15, color=GREEN, align=PP_ALIGN.CENTER)

box(s3, 0.4, 5.7, 12.5, 1.4, DGRAY)
txt(s3, "Principais ameaças cobertas:", 0.7, 5.78, 12.0, 0.45, size=14, bold=True, color=WHITE)
threats = "CVEs em containers  ·  Dependências vulneráveis (SCA)  ·  SQL Injection  ·  XSS  ·  Headers inseguros  ·  CORS  ·  Secrets expostos"
txt(s3, threats, 0.7, 6.2, 12.0, 0.6, size=13, color=LGRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 4 — Arquitetura
# ═══════════════════════════════════════════════════════════════════════════════
s4 = slide()
header(s4, "Arquitetura dos Serviços", "Como os componentes se comunicam", TEAL)

# Browser
box(s4, 0.2, 3.0, 1.8, 1.0, DGRAY)
txt(s4, "🌐\nBrowser", 0.2, 3.0, 1.8, 1.0, size=12, align=PP_ALIGN.CENTER)

txt(s4, "→", 2.05, 3.35, 0.4, 0.4, size=18, color=LGRAY)

# Frontend
box(s4, 2.5, 2.1, 2.3, 2.2, RGBColor(0x00, 0x44, 0x2A))
txt(s4, "🟩 Frontend :80\nnginx + Vue.js", 2.55, 2.5, 2.2, 1.4,
    size=12, color=WHITE, align=PP_ALIGN.CENTER)

txt(s4, "→\n/api/", 4.9, 3.0, 0.7, 0.7, size=11, color=LGRAY, align=PP_ALIGN.CENTER)

# Backend
box(s4, 5.65, 2.1, 2.3, 2.2, RGBColor(0x00, 0x2A, 0x66))
txt(s4, "🟦 Backend\nNestJS :3000\n+ :9464 metrics", 5.7, 2.25, 2.2, 1.8,
    size=11, color=WHITE, align=PP_ALIGN.CENTER)

# Jaeger
box(s4, 8.3, 1.3, 2.4, 1.1, RGBColor(0x66, 0x2A, 0x00))
txt(s4, "🟠 Jaeger :16686", 8.35, 1.55, 2.3, 0.55,
    size=12, color=WHITE, align=PP_ALIGN.CENTER)

# Prometheus
box(s4, 8.3, 2.7, 2.4, 1.1, RGBColor(0x55, 0x18, 0x00))
txt(s4, "🔴 Prometheus :9090", 8.35, 2.95, 2.3, 0.55,
    size=11, color=WHITE, align=PP_ALIGN.CENTER)

txt(s4, "→ OTel gRPC", 8.0, 2.3, 1.0, 0.4, size=9, color=LGRAY)
txt(s4, "→ scrape", 8.0, 3.1, 0.9, 0.35, size=9, color=LGRAY)

# DefectDojo
box(s4, 8.3, 4.2, 2.9, 1.2, RGBColor(0x33, 0x00, 0x66))
txt(s4, "🟣 DefectDojo :8081\nnginx + Django + Celery", 8.35, 4.35, 2.8, 0.9,
    size=11, color=WHITE, align=PP_ALIGN.CENTER)

# AppSec tools area
box(s4, 0.2, 5.3, 7.8, 1.7, DGRAY)
txt(s4, "🛡️  AppSec Scans", 0.4, 5.38, 3.5, 0.5, size=14, bold=True, color=RED)

for tlabel, tx in [("🐳 Trivy", 0.4), ("📦 Dep. Check", 2.3), ("⚡ OWASP ZAP", 4.2)]:
    box(s4, tx, 5.85, 1.8, 0.75, RGBColor(0x3D, 0x1A, 0x1A))
    txt(s4, tlabel, tx + 0.05, 5.9, 1.7, 0.55, size=12, color=WHITE,
        align=PP_ALIGN.CENTER)

txt(s4, "→ reports/ → import →", 6.3, 5.85, 2.2, 0.6, size=11, color=LGRAY)
txt(s4, "↑", 9.7, 4.95, 0.4, 0.35, size=18, color=PURPLE, align=PP_ALIGN.CENTER)

# OTel traces from browser
txt(s4, "OTel\n/otel/ proxy →", 2.55, 4.55, 2.5, 0.7, size=10, color=LGRAY,
    italic=True)
txt(s4, "↗ OTel HTTP", 4.95, 1.65, 1.5, 0.4, size=9, color=LGRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 5 — Backend
# ═══════════════════════════════════════════════════════════════════════════════
s5 = slide()
header(s5, "Etapa 1 — Backend NestJS", "API REST com Swagger, OpenTelemetry e Prometheus", BLUE)

box(s5, 0.4, 1.35, 5.8, 5.8, BG2)
txt(s5, "⚙️  Stack", 0.65, 1.45, 5.3, 0.5, size=18, bold=True, color=BLUE)

for i, (icon, key, val) in enumerate([
    ("🟢", "Runtime",   "Node.js 22"),
    ("🔴", "Framework", "NestJS 11"),
    ("🔵", "Linguagem", "TypeScript 5"),
    ("📄", "Docs",      "Swagger / OpenAPI"),
    ("✅", "Validação", "class-validator + class-transformer"),
    ("📡", "Traces",    "OpenTelemetry SDK → Jaeger"),
    ("📊", "Métricas",  "Prometheus / prom-client"),
    ("💚", "Health",    "@nestjs/terminus"),
]):
    y = 2.05 + i * 0.54
    txt(s5, f"{icon}  {key}:", 0.65, y, 2.6, 0.46, size=13, color=LGRAY)
    txt(s5, val, 3.15, y, 2.8, 0.46, size=13, bold=True)

box(s5, 6.5, 1.35, 6.5, 5.8, BG2)
txt(s5, "🌐  Endpoints", 6.75, 1.45, 5.8, 0.5, size=18, bold=True, color=BLUE)

for i, (method, path, color, desc) in enumerate([
    ("POST", "/bmi/calculate", GREEN,  "Calcula o IMC"),
    ("GET",  "/health",        TEAL,   "Health check da aplicação"),
    ("GET",  "/api/docs",      ORANGE, "Swagger UI"),
    ("GET",  ":9464/metrics",  RED,    "Métricas Prometheus (porta separada)"),
]):
    y = 2.1 + i * 1.25
    pill(s5, method, 6.7, y + 0.08, color, size=12)
    txt(s5, path, 8.1, y + 0.06, 4.6, 0.45, size=14, bold=True)
    txt(s5, desc, 8.1, y + 0.55, 4.6, 0.42, size=13, color=LGRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 6 — Classificações IMC
# ═══════════════════════════════════════════════════════════════════════════════
s6 = slide()
header(s6, "Classificações de IMC (OMS)", "Lógica de negócio do BmiService", BLUE)

for i, (bmi_range, label, color) in enumerate([
    ("< 16",        "Severe Underweight", RGBColor(0x8B, 0x00, 0x00)),
    ("16 – 18.4",   "Underweight",        RED),
    ("18.5 – 24.9", "Normal weight",      GREEN),
    ("25 – 29.9",   "Overweight",         ORANGE),
    ("30 – 34.9",   "Obesity Class I",    RGBColor(0xFF, 0x66, 0x00)),
    ("35 – 39.9",   "Obesity Class II",   RGBColor(0xCC, 0x33, 0x00)),
    ("≥ 40",        "Obesity Class III",  RGBColor(0x99, 0x00, 0x00)),
]):
    y = 1.38 + i * 0.84
    box(s6, 0.4, y, 0.32, 0.74, color)
    box(s6, 0.82, y, 11.9, 0.74, BG2)
    txt(s6, bmi_range, 1.02, y + 0.14, 2.8, 0.46, size=16, color=LGRAY)
    txt(s6, label, 4.2, y + 0.14, 8.3, 0.46, size=18, bold=True)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 7 — Frontend
# ═══════════════════════════════════════════════════════════════════════════════
s7 = slide()
header(s7, "Etapa 2 — Frontend Vue.js", "SPA com OpenTelemetry Web", GREEN)

box(s7, 0.4, 1.35, 5.8, 3.8, BG2)
txt(s7, "🖥️  Stack", 0.65, 1.45, 5.3, 0.5, size=18, bold=True, color=GREEN)

for i, (icon, key, val) in enumerate([
    ("🟩", "Framework", "Vue.js 3 (Composition API)"),
    ("⚡", "Build",     "Vite 6"),
    ("🔵", "Linguagem", "TypeScript 5"),
    ("📡", "HTTP",      "Axios"),
    ("🛤️", "Routing",  "Vue Router"),
    ("📊", "Traces",    "OpenTelemetry Web SDK"),
]):
    y = 2.05 + i * 0.54
    txt(s7, f"{icon}  {key}:", 0.65, y, 2.6, 0.46, size=13, color=LGRAY)
    txt(s7, val, 3.15, y, 3.0, 0.46, size=13, bold=True)

box(s7, 6.5, 1.35, 6.5, 3.8, BG2)
txt(s7, "🔀  Nginx Proxy", 6.75, 1.45, 5.8, 0.5, size=18, bold=True, color=GREEN)

for i, (route, dest, desc, color) in enumerate([
    ("/",      "Vue.js SPA",   "Arquivos estáticos",         WHITE),
    ("/api/",  "backend:3000", "Proxy da API REST",          BLUE),
    ("/otel/", "jaeger:4318",  "Traces OTel (evita CORS)",   ORANGE),
]):
    y = 2.1 + i * 1.1
    pill(s7, route, 6.7, y + 0.05, DGRAY, fg=color, size=13)
    txt(s7, f"→  {dest}", 8.3, y + 0.04, 4.4, 0.42, size=14, bold=True, color=color)
    txt(s7, desc, 8.3, y + 0.5, 4.4, 0.4, size=12, color=LGRAY)

box(s7, 0.4, 5.35, 12.5, 1.8, DGRAY)
txt(s7, "📡  Fluxo de Traces (Browser → Jaeger)", 0.65, 5.42, 10.0, 0.5,
    size=16, bold=True, color=ORANGE)
txt(s7, "Vue.js  →  FetchInstrumentation  →  /otel/v1/traces  →  nginx proxy  →  Jaeger :4318",
    0.65, 5.92, 12.0, 0.42, size=13, color=WHITE)
txt(s7, "Traces correlacionados entre frontend e backend via W3C TraceContext",
    0.65, 6.35, 12.0, 0.42, size=13, color=LGRAY, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 8 — Testes
# ═══════════════════════════════════════════════════════════════════════════════
s8 = slide()
header(s8, "Etapa 3 — Testes", "Jest (backend) + Vitest (frontend)", ORANGE)

box(s8, 0.4, 1.35, 6.0, 5.8, BG2)
txt(s8, "⚙️  Backend — Jest", 0.65, 1.45, 5.5, 0.5, size=18, bold=True, color=BLUE)

for i, (fname, count, desc) in enumerate([
    ("bmi.service.spec.ts",    "9 testes", "Todas 7 classificações + valor IMC + campos"),
    ("bmi.controller.spec.ts", "3 testes", "Definição, resultado correto, delegação"),
    ("app.e2e-spec.ts",        "6 testes", "POST /bmi + 4 validações + GET /health"),
]):
    y = 2.1 + i * 1.55
    box(s8, 0.6, y, 5.6, 1.35, DGRAY)
    txt(s8, fname, 0.78, y + 0.06, 5.2, 0.45, size=14, bold=True, color=BLUE)
    pill(s8, count, 0.78, y + 0.6, GREEN, size=12)
    txt(s8, desc, 0.78, y + 0.96, 5.0, 0.35, size=11, color=LGRAY)

box(s8, 6.65, 1.35, 6.3, 5.8, BG2)
txt(s8, "🖥️  Frontend — Vitest", 6.9, 1.45, 5.8, 0.5, size=18, bold=True, color=GREEN)

for i, (fname, desc) in enumerate([
    ("BmiForm.spec.ts",    "Renderização e envio do formulário"),
    ("BmiResult.spec.ts",  "Cores por classificação, evento reset"),
    ("bmi.service.spec.ts","Chamada à API via Axios"),
    ("useBmi.spec.ts",     "Estados: loading / error / result"),
]):
    y = 2.1 + i * 1.22
    box(s8, 6.8, y, 6.0, 1.05, DGRAY)
    txt(s8, fname, 6.98, y + 0.06, 5.7, 0.42, size=14, bold=True, color=GREEN)
    txt(s8, desc,  6.98, y + 0.58, 5.7, 0.38, size=12, color=LGRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 9 — AppSec Overview
# ═══════════════════════════════════════════════════════════════════════════════
s9 = slide()
header(s9, "Etapa 4 — AppSec", "Ferramentas de segurança integradas ao pipeline", RED)

for i, (icon, name, type_, desc, color) in enumerate([
    ("🐳", "Trivy",       "Container Scan", "CVEs nas imagens Docker\nappsec-bmi-backend/frontend", BLUE),
    ("📦", "Dep. Check",  "SCA",            "Vulnerabilidades em\ndependências npm (NVD)", ORANGE),
    ("⚡", "OWASP ZAP",  "DAST",           "Vulnerabilidades HTTP\nem runtime (full scan)", RED),
    ("🟣", "DefectDojo",  "Aggregator",     "Centraliza, prioriza e\ngerencia todos os findings", PURPLE),
]):
    x = 0.35 + i * 3.25
    box(s9, x, 1.45, 3.05, 5.4, BG2)
    box(s9, x, 1.45, 3.05, 0.06, color)
    txt(s9, icon,  x, 1.55, 3.05, 1.2,  size=52, align=PP_ALIGN.CENTER)
    txt(s9, name,  x, 2.75, 3.05, 0.55, size=20, bold=True, align=PP_ALIGN.CENTER)
    pill(s9, type_, x + 0.55, 3.38, color, size=12)
    txt(s9, desc,  x + 0.1, 3.9, 2.85, 1.5, size=13, color=LGRAY,
        align=PP_ALIGN.CENTER)

    if i < 3:
        txt(s9, "→", x + 3.0, 3.9, 0.3, 0.45, size=20, color=LGRAY,
            align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 10 — Trivy
# ═══════════════════════════════════════════════════════════════════════════════
s10 = slide()
header(s10, "Trivy — Container Scan", "Identificação de CVEs nas imagens Docker", BLUE)

box(s10, 0.4, 1.35, 5.8, 5.8, BG2)
txt(s10, "🐳  O que analisa?", 0.65, 1.45, 5.3, 0.5, size=18, bold=True, color=BLUE)

for i, item in enumerate([
    "📦  Pacotes do OS (Alpine / Debian)",
    "🔷  Dependências da aplicação",
    "🔑  Secrets e credenciais expostas",
    "⚙️  Configurações inseguras",
    "📋  Licenças de software",
]):
    txt(s10, item, 0.65, 2.1 + i * 0.65, 5.3, 0.55, size=15, color=WHITE)

box(s10, 0.4, 5.4, 5.8, 1.7, BG2)
txt(s10, "🎯  Severidades:", 0.65, 5.5, 5.3, 0.42, size=14, bold=True, color=LGRAY)
for sev, color, sx in [
    ("CRITICAL", RGBColor(0x99, 0, 0), 0.65),
    ("HIGH",     RED,                   2.5),
    ("MEDIUM",   ORANGE,                4.1),
]:
    pill(s10, sev, sx, 6.0, color, size=12)

box(s10, 6.5, 1.35, 6.5, 2.8, RGBColor(0x08, 0x08, 0x08))
for i, line in enumerate([
    "$ trivy image \\",
    "  --format json \\",
    "  --severity CRITICAL,HIGH,MEDIUM \\",
    "  appsec-bmi-backend:latest",
]):
    color = GREEN if i == 0 else (YELLOW if i == 3 else WHITE)
    txt(s10, line, 6.65, 1.45 + i * 0.55, 6.15, 0.48, size=13, color=color)

box(s10, 6.5, 4.4, 6.5, 2.7, BG2)
txt(s10, "📊  Resultados da última execução", 6.75, 4.5, 6.0, 0.5,
    size=16, bold=True)
txt(s10, "🔴  Backend:   3 vulnerabilidades",  6.75, 5.1,  6.0, 0.45, size=15, color=RED)
txt(s10, "🟠  Frontend: 77 vulnerabilidades",  6.75, 5.58, 6.0, 0.45, size=15, color=ORANGE)
txt(s10, "(maioria: pacotes OS Alpine 3.21)",   6.75, 6.05, 6.0, 0.4,  size=13,
    color=LGRAY, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 11 — OWASP Dependency Check
# ═══════════════════════════════════════════════════════════════════════════════
s11 = slide()
header(s11, "OWASP Dependency Check — SCA",
       "Software Composition Analysis das dependências npm", ORANGE)

box(s11, 0.4, 1.35, 5.8, 5.8, BG2)
txt(s11, "📦  O que analisa?", 0.65, 1.45, 5.3, 0.5, size=18, bold=True, color=ORANGE)

for i, item in enumerate([
    "📄  package.json / package-lock.json",
    "🗂️  node_modules (backend + frontend)",
    "🔗  CVEs no banco NVD — NIST (346k+)",
    "📋  Vulnerabilidades por versão de pacote",
    "🏷️  Licenças permissivas / restritivas",
]):
    txt(s11, item, 0.65, 2.1 + i * 0.65, 5.3, 0.55, size=14, color=WHITE)

box(s11, 6.5, 1.35, 6.5, 2.5, RGBColor(0x08, 0x08, 0x08))
for i, line in enumerate([
    "$ docker run owasp/dependency-check \\",
    "  --scan /src/backend \\",
    "  --scan /src/frontend \\",
    "  --format JSON --format HTML",
]):
    color = GREEN if i == 0 else (YELLOW if i == 3 else WHITE)
    txt(s11, line, 6.65, 1.45 + i * 0.5, 6.15, 0.45, size=12, color=color)

box(s11, 6.5, 4.1, 6.5, 3.0, BG2)
txt(s11, "⚠️  Sobre o NVD Database", 6.75, 4.2, 6.0, 0.5, size=16, bold=True, color=YELLOW)
for i, (item, color) in enumerate([
    ("• Banco com 346.000+ CVEs registrados",      WHITE),
    ("• Primeira execução: ~30 min sem API key",   ORANGE),
    ("• Volume Docker persiste — próximas = rápido",GREEN),
    ("• NVD API Key gratuita → 10× mais rápido",   LGRAY),
]):
    txt(s11, item, 6.75, 4.8 + i * 0.52, 6.0, 0.44, size=13, color=color)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 12 — OWASP ZAP
# ═══════════════════════════════════════════════════════════════════════════════
s12 = slide()
header(s12, "OWASP ZAP — DAST",
       "Dynamic Application Security Testing em runtime", RED)

box(s12, 0.4, 1.35, 5.8, 5.8, BG2)
txt(s12, "⚡  O que analisa?", 0.65, 1.45, 5.3, 0.5, size=18, bold=True, color=RED)

for i, item in enumerate([
    "🕷️  Crawler automático (full scan)",
    "💉  SQL Injection / XSS",
    "🔓  Falhas de autenticação",
    "📡  Headers HTTP inseguros",
    "🔒  Configuração de CORS",
    "🗺️  Endpoints expostos indevidamente",
    "🔁  CSRF e clickjacking",
]):
    txt(s12, item, 0.65, 2.1 + i * 0.6, 5.3, 0.52, size=14, color=WHITE)

box(s12, 6.5, 1.35, 6.5, 2.0, BG2)
txt(s12, "⚙️  Regras suprimidas (zap-rules.tsv)", 6.75, 1.45, 6.0, 0.5,
    size=14, bold=True, color=WHITE)
for i, rule in enumerate([
    "IGNORE  10015  Cache-control Header",
    "IGNORE  10096  Timestamp Disclosure",
    "IGNORE  10027  Suspicious Comments",
]):
    txt(s12, rule, 6.75, 2.0 + i * 0.42, 6.1, 0.38, size=13, color=LGRAY)

box(s12, 6.5, 3.6, 6.5, 3.5, BG2)
txt(s12, "🔄  Fluxo de execução", 6.75, 3.7, 6.0, 0.5, size=16, bold=True, color=RED)
for i, (num, step, color) in enumerate([
    ("1", "ZAP acessa http://localhost:80",        WHITE),
    ("2", "Crawler mapeia todos os endpoints",     WHITE),
    ("3", "Executa testes de vulnerabilidade",     WHITE),
    ("4", "Gera JSON + HTML em reports/",          WHITE),
    ("5", "Importa findings no DefectDojo",        GREEN),
]):
    pill(s12, num, 6.65, 4.3 + i * 0.5, DGRAY, size=11)
    txt(s12, step, 7.15, 4.3 + i * 0.5, 5.6, 0.42, size=13, color=color)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 13 — DefectDojo
# ═══════════════════════════════════════════════════════════════════════════════
s13 = slide()
header(s13, "DefectDojo — Aggregator",
       "Centralização e gestão de vulnerabilidades", PURPLE)

box(s13, 0.4, 1.35, 5.8, 5.8, BG2)
txt(s13, "🟣  O que é?", 0.65, 1.45, 5.3, 0.5, size=18, bold=True, color=PURPLE)

for i, item in enumerate([
    "📋  Agrega findings de todas as ferramentas",
    "🎯  Prioriza por severidade e risco",
    "✅  Controle de falso-positivos",
    "📈  Histórico de vulnerabilidades",
    "🔔  Deduplication automática",
    "🔗  Integração com Jira, Slack e mais",
]):
    txt(s13, item, 0.65, 2.1 + i * 0.65, 5.3, 0.55, size=14, color=WHITE)

box(s13, 6.5, 1.35, 6.5, 5.8, BG2)
txt(s13, "🐳  Stack de Containers", 6.75, 1.45, 6.0, 0.5, size=18, bold=True, color=PURPLE)

for i, (name, desc, color) in enumerate([
    ("defectdojo-nginx",       "Proxy + assets estáticos",         PURPLE),
    ("uwsgi",                  "Django app via uWSGI",             PURPLE),
    ("defectdojo-celeryworker","Processa imports em background",    BLUE),
    ("defectdojo-celerybeat",  "Agendamento de tarefas periódicas", BLUE),
    ("redis",                  "Broker de tarefas Celery",          RED),
    ("postgres",               "Banco de dados principal",          TEAL),
]):
    y = 2.1 + i * 0.85
    box(s13, 6.6, y, 0.07, 0.72, color)
    txt(s13, name, 6.85, y + 0.03, 3.5, 0.36, size=13, bold=True)
    txt(s13, desc, 6.85, y + 0.42, 6.1, 0.3,  size=11, color=LGRAY)

box(s13, 0.4, 6.5, 12.5, 0.65, DGRAY)
txt(s13, "🔗  http://localhost:8081  •  admin / admin@dojo123",
    0.7, 6.6, 12.0, 0.42, size=14, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 14 — GitHub Actions
# ═══════════════════════════════════════════════════════════════════════════════
s14 = slide()
header(s14, "GitHub Actions — CI/CD Pipeline",
       "Scans automáticos em cada push / PR para main", TEAL)

for i, (icon, job, tool, desc, color) in enumerate([
    ("🧪", "test",             "Jest + Vitest",    "Valida que nenhum teste quebrou",     GREEN),
    ("📦", "dependency-check", "OWASP Dep. Check", "SCA das dependências npm",            ORANGE),
    ("🐳", "trivy",            "Trivy",            "Scan das imagens backend e frontend", BLUE),
    ("⚡", "zap",              "OWASP ZAP",        "DAST contra o stack completo",        RED),
]):
    x = 0.35 + i * 3.25
    box(s14, x, 1.45, 3.05, 0.42, color)
    txt(s14, f"{icon}  {job}", x + 0.1, 1.47, 2.85, 0.38, size=13, bold=True)
    box(s14, x, 1.87, 3.05, 3.3, BG2)
    txt(s14, tool, x + 0.1, 1.97, 2.85, 0.55, size=17, bold=True)
    txt(s14, desc, x + 0.1, 2.62, 2.85, 1.3, size=13, color=LGRAY)
    if i < 3:
        txt(s14, "→", x + 3.0, 3.1, 0.3, 0.45, size=18, color=LGRAY,
            align=PP_ALIGN.CENTER)

box(s14, 0.4, 5.4, 12.5, 0.75, BG2)
txt(s14, "📁  Artefatos:", 0.65, 5.5, 2.0, 0.42, size=13, bold=True)
txt(s14, "trivy-report.json  ·  dependency-check-report.html  ·  zap-report.json  ·  zap-report.html",
    2.5, 5.5, 10.2, 0.42, size=13, color=LGRAY)

box(s14, 0.4, 6.3, 12.5, 0.85, DGRAY)
txt(s14, "🟣  Secrets:  DEFECTDOJO_URL  +  DEFECTDOJO_TOKEN  +  DEFECTDOJO_PRODUCT_ID  →  import automático",
    0.65, 6.45, 12.0, 0.45, size=13, color=PURPLE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 15 — Princípios Aplicados
# ═══════════════════════════════════════════════════════════════════════════════
s15 = slide()
header(s15, "Princípios Aplicados", "Boas práticas de desenvolvimento seguro", TEAL)

for i, (icon, title, desc, color) in enumerate([
    ("🏗️", "SOLID / MVC",          "Controller → Service → DTO — responsabilidade única",               BLUE),
    ("✅",  "Validação na borda",   "ValidationPipe global: whitelist + transform + rejeita extras",     GREEN),
    ("🔒",  "Container seguro",     "node:22-alpine, usuário não-root, multi-stage build",               ORANGE),
    ("📊",  "Métricas isoladas",    ":9464 não exposta publicamente — apenas para o Prometheus interno", TEAL),
    ("📡",  "Observabilidade",      "Traces correlacionados via W3C TraceContext (frontend ↔ backend)",  PURPLE),
    ("🛡️", "Segurança contínua",   "SCA + Container Scan + DAST automatizados em cada push",            RED),
]):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.55
    y = 1.45 + row * 1.9
    box(s15, x, y, 0.06, 1.6, color)
    txt(s15, f"{icon}  {title}", x + 0.2, y + 0.08, 6.0, 0.52, size=17, bold=True)
    txt(s15, desc, x + 0.2, y + 0.68, 6.0, 0.72, size=13, color=LGRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 16 — Conclusão
# ═══════════════════════════════════════════════════════════════════════════════
s16 = slide()
box(s16, 0, 0, 13.33, 0.07, GREEN)

txt(s16, "🎯", 3.5, 0.4, 6.33, 1.3, size=90, align=PP_ALIGN.CENTER)
txt(s16, "Conclusão", 0.5, 1.6, 12.33, 0.95,
    size=46, bold=True, align=PP_ALIGN.CENTER)
txt(s16, "AppSec não é um produto — é uma cultura",
    0.5, 2.58, 12.33, 0.65, size=22, color=LGRAY, italic=True, align=PP_ALIGN.CENTER)

box(s16, 2.0, 3.35, 9.33, 3.3, BG2)
for i, item in enumerate([
    "✅  Segurança integrada ao pipeline, não uma etapa final",
    "✅  Ferramentas open-source e gratuitas já resolvem muito",
    "✅  Automação via GitHub Actions — zero esforço por deploy",
    "✅  DefectDojo centraliza e prioriza o que realmente importa",
    "✅  Shift Left = custo menor, qualidade maior",
]):
    txt(s16, item, 2.3, 3.5 + i * 0.58, 8.8, 0.5, size=15, color=WHITE)

box(s16, 0, 7.1, 13.33, 0.4, GREEN)
txt(s16, "🙋  Perguntas?  •  github.com/questores/appsec-bmi",
    0.5, 7.12, 12.33, 0.35, size=15, bold=True, align=PP_ALIGN.CENTER)


# ─── Save ─────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "..",
                   "appsec-palestra.pptx")
out = os.path.normpath(out)
prs.save(out)
print(f"Salvo: {out}")
